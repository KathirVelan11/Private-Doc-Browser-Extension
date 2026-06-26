"""Document loaders.

Each loader returns a list of ``(text, page)`` units so that page numbers
survive into the retrieval results. A "page" is logical: PDFs and slide
decks have real pages/slides, while txt/md/docx/html/csv are treated as a
single page-1 unit.
"""
import csv as _csv
import os
from html.parser import HTMLParser
from typing import List, Tuple

Unit = Tuple[str, int]  # (text, page_number)


def _load_pdf(path: str) -> List[Unit]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    units: List[Unit] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            units.append((text, i))
    return units


def _load_docx(path: str) -> List[Unit]:
    import docx

    document = docx.Document(path)
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [(text, 1)] if text.strip() else []


def _load_pptx(path: str) -> List[Unit]:
    from pptx import Presentation

    prs = Presentation(path)
    units: List[Unit] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        text = "\n".join(p for p in parts if p.strip())
        if text.strip():
            units.append((text, i))
    return units


class _TextExtractor(HTMLParser):
    """Strip tags, drop <script>/<style> content, keep readable text."""

    def __init__(self):
        super().__init__()
        self._skip = 0
        self.chunks: List[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip:
            self._skip -= 1

    def handle_data(self, data):
        if not self._skip and data.strip():
            self.chunks.append(data.strip())


def _load_html(path: str) -> List[Unit]:
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        parser = _TextExtractor()
        parser.feed(fh.read())
    text = "\n".join(parser.chunks)
    return [(text, 1)] if text.strip() else []


def _load_csv(path: str) -> List[Unit]:
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as fh:
        for row in _csv.reader(fh):
            cells = [c.strip() for c in row if c.strip()]
            if cells:
                rows.append(" | ".join(cells))
    text = "\n".join(rows)
    return [(text, 1)] if text.strip() else []


def _load_text(path: str) -> List[Unit]:
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        text = fh.read()
    return [(text, 1)] if text.strip() else []


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".html": _load_html,
    ".htm": _load_html,
    ".csv": _load_csv,
    ".txt": _load_text,
    ".md": _load_text,
}


def load_document(path: str) -> List[Unit]:
    """Load a single document into a list of (text, page) units."""
    ext = os.path.splitext(path)[1].lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return loader(path)


def is_supported(filename: str) -> bool:
    return os.path.splitext(filename)[1].lower() in _LOADERS
